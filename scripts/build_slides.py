"""
Build a PowerPoint deck from a slide outline.

Usage:
    uv run python scripts/build_slides.py \
        slides/term01/lec07_surface_consistent_deconvolution/slide_outline.md \
        -o slides/term01/lec07_surface_consistent_deconvolution/lec07_surface_consistent_deconvolution.pptx
"""

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def latex_to_plain(text: str) -> str:
    """Convert simple LaTeX math to plain text for PowerPoint."""
    # Remove inline/display math delimiters
    text = re.sub(r"\$\$(.+?)\$\$", lambda m: _convert_math(m.group(1)), text, flags=re.DOTALL)
    text = re.sub(r"\$(.+?)\$", lambda m: _convert_math(m.group(1)), text)
    # Remove stray \,
    text = text.replace(r"\,", " ")
    return text


def _convert_math(expr: str) -> str:
    """Best-effort conversion of a LaTeX math expression to plain text."""
    # Common substitutions
    replacements = {
        r"\varepsilon": "ε",
        r"\epsilon": "ε",
        r"\alpha": "α",
        r"\beta": "β",
        r"\gamma": "γ",
        r"\delta": "δ",
        r"\Delta": "Δ",
        r"\sigma": "σ",
        r"\Sigma": "Σ",
        r"\phi": "φ",
        r"\Phi": "Φ",
        r"\theta": "θ",
        r"\omega": "ω",
        r"\Omega": "Ω",
        r"\lambda": "λ",
        r"\mu": "μ",
        r"\nu": "ν",
        r"\pi": "π",
        r"\tau": "τ",
        r"\cdot": "·",
        r"\times": "×",
        r"\rightarrow": "→",
        r"\leftarrow": "←",
        r"\leftrightarrow": "↔",
        r"\approx": "≈",
        r"\neq": "≠",
        r"\leq": "≤",
        r"\geq": "≥",
        r"\pm": "±",
        r"\infty": "∞",
        r"\partial": "∂",
        r"\nabla": "∇",
        r"\int": "∫",
        r"\sum": "Σ",
        r"\prod": "Π",
        r"\mathrm{}": "",
        r"\text{}": "",
        r"\mathbf{R}": "R",
        r"\mathbf{f}": "f",
        r"\mathbf{d}": "d",
    }
    for latex, plain in replacements.items():
        expr = expr.replace(latex, plain)

    # Remove remaining \mathbf{...}, \mathrm{...}, \text{...} wrappers
    expr = re.sub(r"\\mathbf\{([^}]*)\}", r"\1", expr)
    expr = re.sub(r"\\mathrm\{([^}]*)\}", r"\1", expr)
    expr = re.sub(r"\\text\{([^}]*)\}", r"\1", expr)

    # Superscripts: x^2, x^{abc}
    expr = re.sub(r"\^\{([^}]*)\}", lambda m: _to_superscript(m.group(1)), expr)
    expr = re.sub(r"\^([0-9a-zA-Z])", lambda m: _to_superscript(m.group(1)), expr)

    # Subscripts: x_2, x_{abc}
    expr = re.sub(r"_\{([^}]*)\}", lambda m: _to_subscript(m.group(1)), expr)
    expr = re.sub(r"_([0-9a-zA-Z])", lambda m: _to_subscript(m.group(1)), expr)

    # Fractions: \frac{a}{b} -> a/b
    expr = re.sub(r"\\frac\{([^}]*)\}\{([^}]*)\}", r"(\1)/(\2)", expr)

    # Remove remaining backslashes (e.g. \max, \ln) — keep the command name
    expr = re.sub(r"\\([a-zA-Z]+)", r"\1", expr)

    # Clean up braces
    expr = expr.replace("{", "").replace("}", "")
    # Replace \* leftover
    expr = expr.replace("\\", "")

    return expr.strip()


def _to_superscript(s: str) -> str:
    mapping = str.maketrans("0123456789+-=()abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ", "⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁼⁽⁾ᵃᵇᶜᵈᵉᶠᵍʰⁱʲᵏˡᵐⁿᵒᵖᵠʳˢᵗᵘᵛʷˣʸᶻᴬᴮᶜᴰᴱᶠᴳᴴᴵᴶᴷᴸᴹᴺᴼᴾᵠᴿˢᵀᵁⱽᵂˣʸᶻ")
    return s.translate(mapping)


def _to_subscript(s: str) -> str:
    mapping = str.maketrans("0123456789+-=()abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ", "₀₁₂₃₄₅₆₇₈₉₊₋₌₍₎ₐբ꜀ᴅₑբɢʜᵢᴊₖₗₘₙₒᵽᵩʀₛᵗᵤᵥᴡₓʏᴢₐʙᴄᴅᴇꜰɢʜɪᴊᴋʟᴍɴᴏᴘǫʀѕᴛᴜᴠᴡхʏᴢ")
    return s.translate(mapping)


def strip_backticks(text: str) -> str:
    """Remove inline code backticks."""
    return re.sub(r"`([^`]+)`", r"\1", text)


def parse_outline(text: str):
    """Parse a markdown slide outline into a list of slide dicts.

    Supports two slide formats:
      - Sections separated by `---` with a leading `# Title` line.
      - Explicit `## Slide N — Title` markers.
    """
    slides = []
    current = None
    started = False  # becomes True after the first `---` separator

    for line in text.splitlines():
        line = line.rstrip()
        if not line:
            continue

        # Section separator: reset for the next slide
        if line.strip() == "---":
            started = True
            if current:
                slides.append(current)
                current = None
            continue

        # Skip the document title before the first separator
        if not started and line.startswith("#"):
            continue

        # Explicit slide marker: ## Slide N — Title
        m = re.match(r"^##\s+Slide\s+\d+\s*—\s*(.+)$", line)
        if m:
            if current:
                slides.append(current)
            current = {"title": m.group(1).strip(), "bullets": [], "figure": None}
            continue

        # Slide title: # Title
        if started and line.startswith("# "):
            if current:
                slides.append(current)
            current = {"title": line[2:].strip(), "bullets": [], "figure": None}
            continue

        if current is None:
            continue

        # Figure reference (with optional backticks and optional bold markdown)
        fig_m = re.match(r"^(?:[-*]\s+|\*\*)?Figure:\*\*?\s*`?([^`]+)`?$", line, re.IGNORECASE)
        if fig_m:
            current["figure"] = fig_m.group(1).strip()
            if "optional" in current["figure"].lower() or "can be added" in current["figure"].lower():
                current["figure"] = None
            continue

        # Bullet point
        bullet_m = re.match(r"^[-*]\s+(.*)$", line)
        if bullet_m:
            txt = strip_backticks(bullet_m.group(1).strip())
            txt = latex_to_plain(txt)
            current["bullets"].append(txt)
            continue

        # Sub-bullet (indent with spaces then - or *)
        sub_m = re.match(r"^\s+[-*]\s+(.*)$", line)
        if sub_m and current["bullets"]:
            txt = strip_backticks(sub_m.group(1).strip())
            txt = latex_to_plain(txt)
            current["bullets"][-1] += "\n    • " + txt
            continue

        # Plain text that is not a bullet: treat as a continuation bullet
        plain = strip_backticks(line.strip())
        plain = latex_to_plain(plain)
        if current["bullets"]:
            current["bullets"][-1] += " " + plain
        else:
            current["bullets"].append(plain)

    if current:
        slides.append(current)

    return slides


def add_title_slide(prs, title: str, subtitle: str = ""):
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title: str, bullets: list[str], figure_path: str | None, out_dir: Path):
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(10)

    if figure_path:
        fig = Path(figure_path)
        path = None
        if fig.is_absolute() and fig.exists():
            path = str(fig)
        elif fig.exists():
            path = str(fig)
        else:
            # Lecture number prefix heuristic: term01_lec06_xxx.png, term01_lec07_xxx.png, etc.
            lec_prefix = None
            m = re.match(r"(term\d{2}_lec\d{2})_", fig.name)
            if m:
                lec_prefix = m.group(1)
            candidates = [
                out_dir / fig,
                Path("figures") / fig,
                Path("figures") / lec_prefix / fig if lec_prefix else None,
                out_dir.parent.parent.parent / "figures" / fig,
                out_dir.parent.parent.parent / "figures" / lec_prefix / fig if lec_prefix else None,
            ]
            candidates = [c for c in candidates if c is not None]
            for candidate in candidates:
                if candidate.exists():
                    path = str(candidate)
                    break

        if path and Path(path).exists():
            # Place figure on the right half of the slide
            left = Inches(6.8)
            top = Inches(1.6)
            width = Inches(6.0)
            height = Inches(5.5)
            slide.shapes.add_picture(path, left, top, width=width, height=height)
            # Shrink text box to make room
            body.left = Inches(0.5)
            body.top = Inches(1.6)
            body.width = Inches(6.0)
            body.height = Inches(5.5)
        else:
            print(f"  Warning: figure not found: {figure_path}")

    return slide


def main():
    parser = argparse.ArgumentParser(description="Build PowerPoint from slide outline")
    parser.add_argument("outline", help="Path to markdown slide outline")
    parser.add_argument("-o", "--output", required=True, help="Output PPTX path")
    args = parser.parse_args()

    outline_path = Path(args.outline)
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    text = outline_path.read_text(encoding="utf-8")
    slides = parse_outline(text)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    first = slides[0] if slides else {"title": "", "bullets": [], "figure": None}
    title = first["title"]
    subtitle = ""
    # If the first slide has bullets, use them as subtitle lines
    if first["bullets"]:
        subtitle = "\n".join(first["bullets"])
    add_title_slide(prs, title, subtitle)

    for slide_info in slides[1:]:
        add_content_slide(prs, slide_info["title"], slide_info["bullets"], slide_info["figure"], output_path.parent)

    prs.save(str(output_path))
    print(f"Saved {output_path}")


if __name__ == "__main__":
    main()
