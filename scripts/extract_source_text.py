#!/usr/bin/env python3
"""
Extract text from PDF and PPT/PPTX files for LLM ingestion.

Usage:
    python scripts/extract_source_text.py <input_file_or_dir> [--output <dir>] [--pages <range>]

Examples:
    # Extract a whole PDF
    python scripts/extract_source_text.py papers/deconvolution/file.pdf

    # Extract only pages 10-20 of a PDF
    python scripts/extract_source_text.py papers/deconvolution/file.pdf --pages 10-20

    # Extract all supported files in a directory
    python scripts/extract_source_text.py papers/deconvolution --output wiki/sources/_raw_text
"""

import argparse
import subprocess
import sys
from pathlib import Path

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def extract_pdf_text(path: Path, page_range: str | None = None) -> str:
    """Extract text from PDF using PyMuPDF if available, else pdftotext."""
    if fitz is not None:
        doc = fitz.open(path)
        pages = []
        start, end = _parse_page_range(page_range, len(doc))
        for i in range(start, end):
            page = doc.load_page(i)
            pages.append(f"--- Page {i + 1} ---\n{page.get_text()}")
        doc.close()
        return "\n\n".join(pages)
    else:
        cmd = ["pdftotext"]
        if page_range:
            start, end = page_range.split("-")
            cmd.extend(["-f", start.strip(), "-l", end.strip()])
        cmd.extend([str(path), "-"])
        return subprocess.run(cmd, stdout=subprocess.PIPE, text=True).stdout


def extract_pptx_text(path: Path) -> str:
    """Extract text and notes from a .pptx file."""
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed")
    prs = Presentation(path)
    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        lines = [f"--- Slide {idx} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
            if hasattr(shape, "notes_slide") and shape.notes_slide:
                notes = shape.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"[Notes]: {notes}")
        slides.append("\n".join(lines))
    return "\n\n".join(slides)


def extract_ppt_text(path: Path) -> str:
    """Fallback for old .ppt files: convert to text via catppt if available, else note limitation."""
    result = subprocess.run(["which", "catppt"], stdout=subprocess.DEVNULL)
    if result.returncode == 0:
        return subprocess.run(["catppt", str(path)], stdout=subprocess.PIPE, text=True).stdout
    raise RuntimeError("Legacy .ppt extraction requires 'catppt'. Install with your package manager.")


def _parse_page_range(page_range: str | None, max_pages: int) -> tuple[int, int]:
    if not page_range:
        return 0, max_pages
    if "-" in page_range:
        start, end = page_range.split("-", 1)
        return max(0, int(start.strip()) - 1), min(max_pages, int(end.strip()))
    page = int(page_range.strip()) - 1
    return max(0, page), min(max_pages, page + 1)


def extract_file(path: Path, page_range: str | None = None) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path, page_range)
    elif suffix == ".pptx":
        return extract_pptx_text(path)
    elif suffix == ".ppt":
        return extract_ppt_text(path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


def main():
    parser = argparse.ArgumentParser(description="Extract text from source documents.")
    parser.add_argument("inputs", nargs="+", help="Input file(s) or directory/ies")
    parser.add_argument("--output", "-o", help="Output directory (default: same as input)")
    parser.add_argument("--pages", help="PDF page range, e.g. '10-20' or '5'")
    args = parser.parse_args()

    output_dir = Path(args.output) if args.output else Path("wiki/sources/_raw_text")
    output_dir.mkdir(parents=True, exist_ok=True)

    files = []
    for inp in args.inputs:
        p = Path(inp)
        if p.is_file():
            files.append(p)
        else:
            files.extend(sorted(q for q in p.iterdir() if q.suffix.lower() in {".pdf", ".ppt", ".pptx"}))

    for f in files:
        try:
            text = extract_file(f, args.pages)
            out_file = output_dir / f"{f.stem}.txt"
            out_file.write_text(text, encoding="utf-8")
            print(f"Extracted: {f.name} -> {out_file}")
        except Exception as e:
            print(f"FAILED: {f.name} - {e}", file=sys.stderr)


if __name__ == "__main__":
    main()
